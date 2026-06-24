from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
PRIMARY = RGBColor(0xDC, 0x35, 0x45)  # 红色主题
SECONDARY = RGBColor(0xF8, 0x9F, 0x1B)  # 暖橙色
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
GREEN = RGBColor(0x19, 0x87, 0x54)
BLUE = RGBColor(0x0D, 0x6E, 0xFD)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=None, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if color:
        p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=18, color=None, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        if color:
            p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return txBox

# ========== Slide 1: Title ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), PRIMARY)

add_text_box(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(1.5),
             "🛒 EasyShop 电商平台", font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(1),
             "期末考核项目扩展 · 五大功能升级", font_size=28, color=SECONDARY, bold=False)
add_text_box(slide, Inches(1.2), Inches(5.0), Inches(8), Inches(0.6),
             "基于 FastAPI + React + MySQL 的全功能电商系统", font_size=18, color=GRAY)
add_text_box(slide, Inches(1.2), Inches(5.6), Inches(8), Inches(0.6),
             "演讲时间：10分钟", font_size=16, color=GRAY)

# ========== Slide 2: Agenda ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.8),
             "📋 演讲大纲", font_size=36, color=WHITE, bold=True)

agenda_items = [
    "1️⃣  项目概览与技术栈",
    "2️⃣  功能一：商品收藏系统",
    "3️⃣  功能二：收货地址管理",
    "4️⃣  功能三：商品评论与评分",
    "5️⃣  功能四：优惠券系统",
    "6️⃣  功能五：订单物流追踪",
    "7️⃣  架构设计与实现亮点",
    "8️⃣  总结与致谢",
]
add_bullet_list(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5),
                agenda_items, font_size=26, color=DARK, spacing=Pt(14))

# ========== Slide 3: Project Overview ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.8),
             "📌 项目概览", font_size=36, color=WHITE, bold=True)

# Tech stack boxes
techs = [
    ("前端", "React 19\nTypeScript\nVite + Bootstrap 5\nZustand", BLUE),
    ("后端", "Python 3.12\nFastAPI\nSQLAlchemy 2.0\nJWT 认证", GREEN),
    ("数据库", "MySQL 8.0\n异步 ORM\n连接池管理", SECONDARY),
]
for i, (title, desc, color) in enumerate(techs):
    x = Inches(1.0 + i * 4.0)
    y = Inches(1.8)
    box = add_shape(slide, x, y, Inches(3.5), Inches(3.5), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.9), Inches(0.6),
                 title, font_size=28, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(2.9), Inches(2.0),
                 desc, font_size=18, color=WHITE)

add_text_box(slide, Inches(1.0), Inches(5.8), Inches(11), Inches(0.8),
             "⭐ 扩展功能：商品收藏 | 地址管理 | 评论评分 | 优惠券系统 | 物流追踪",
             font_size=20, color=DARK, bold=True)

# ========== Slide 4: Favorites ==========
def create_feature_slide(slide_num, icon, title, desc_items, tech_items, color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), color)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(12), Inches(0.8),
                 f"{icon} {title}", font_size=36, color=WHITE, bold=True)

    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
                 "💡 功能说明", font_size=24, color=color, bold=True)
    add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(3.5),
                    desc_items, font_size=20, color=DARK, spacing=Pt(10))

    add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
                 "⚙️ 技术实现", font_size=24, color=color, bold=True)
    add_bullet_list(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(3.5),
                    tech_items, font_size=20, color=DARK, spacing=Pt(10))

    # Bottom decoration
    add_shape(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), LIGHT_BG)
    add_text_box(slide, Inches(0.8), Inches(6.85), Inches(12), Inches(0.5),
                 f"📱 前端展示见演示环节  |  第 {slide_num} / 8 部分",
                 font_size=16, color=GRAY)

    return slide

create_feature_slide(1, "⭐", "功能一：商品收藏（Wishlist）",
    [
        "• 用户一键收藏/取消收藏商品",
        "• 商品详情页心形按钮切换",
        "• 独立收藏夹页面（Grid 展示）",
        "• 支持取消收藏、加入购物车",
        "• 防止重复收藏（数据库唯一约束）",
        "• Header 新增收藏快捷入口",
    ],
    [
        "• Models: Favorite（user_id, product_id, UNIQUE约束）",
        "• CRUD: add/remove/check/is_favorited",
        "• API: POST/DELETE/GET /favorites",
        "• Frontend: FavoritesPage + 收藏按钮组件",
        "• 使用 selectinload 预加载商品信息",
    ],
    PRIMARY
)

# ========== Slide 5: Address Management ==========
create_feature_slide(2, "📍", "功能二：收货地址管理",
    [
        "• 多地址 CRUD（最多 10 个）",
        "• 默认地址设置/切换",
        "• 地址编辑/删除/设为默认",
        "• 地址表单字段验证（11位手机号）",
        "• 新建地址时自动排序（默认置顶）",
        "• 前端模态框编辑，用户体验流畅",
    ],
    [
        "• Models: Address（省/市/区/详细地址）",
        "• 级联取消默认：设新默认时自动清旧默认",
        "• API: GET/POST/PUT/DELETE /addresses",
        "• 地址数量硬限制：max 10 个",
        "• 前后端双重手机号验证",
    ],
    BLUE
)

# ========== Slide 6: Reviews ==========
create_feature_slide(3, "📝", "功能三：商品评论与评分",
    [
        "• 1-5 星评分系统（半星支持）",
        "• 评分分布统计（柱状图显示）",
        "• 平均评分 + 评论总数展示",
        "• 评论列表（按时间倒序）",
        "• 防重复评论（一人一评）",
        "• 商品详情页评论区集成",
        "• 评分星标可视化（Bootstrap Icons）",
    ],
    [
        "• Models: Review（rating 1-5, content）",
        "• 评分统计：SQL 聚合查询",
        "• GROUP BY rating 计算分布",
        "• AVG + COUNT 统计平均值",
        "• API: GET/POST/PUT/DELETE /reviews",
        "• 商品详情页集成 ReviewsSection 组件",
    ],
    GREEN
)

# ========== Slide 7: Coupons ==========
create_feature_slide(4, "🎟️", "功能四：优惠券系统",
    [
        "• 两种类型：固定金额/百分比折扣",
        "• 管理员后台 CRUD 管理优惠券",
        "• 用户领取优惠券（按券码）",
        "• 我的优惠券列表（未使用/已使用）",
        "• 下单时计算折扣（Apply API）",
        "• 有效期检查 + 使用次数限制",
        "• 最低消费金额校验",
    ],
    [
        "• Models: Coupon + UserCoupon（多对多）",
        "• CouponType 枚举: FIXED / PERCENT",
        "• 管理员 API: /admin/coupons CRUD",
        "• 用户 API: claim/apply/mine",
        "• 折扣计算引擎（金额上限保护）",
        "• 使用次数 + 重复领取双重校验",
    ],
    SECONDARY
)

# ========== Slide 8: Logistics ==========
create_feature_slide(5, "📦", "功能五：订单物流追踪",
    [
        "• 模拟物流轨迹（5个状态流转）",
        "• 时间线 UI 展示（竖排时间轴）",
        "• 支持一键推进物流状态",
        "• 物流信息：单号/快递公司/状态",
        "• 预计送达时间 + 实时状态更新",
        "• 订单详情页「查看物流」入口",
        "• 已发货/已完成订单可跟踪",
    ],
    [
        "• Models: LogisticsTracking（状态枚举）",
        "• LogisticsStatus: pending→delivered 六阶段",
        "• 事件 JSON 存储（灵活扩展）",
        "• 模拟推进 API：按序流转",
        "• 每次推进自动生成中文事件文案",
        "• Timeline 组件 + 状态着色",
    ],
    RGBColor(0x66, 0x33, 0x99)
)

# ========== Slide 9: Architecture ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK)
add_text_box(slide, Inches(0.8), Inches(0.2), Inches(12), Inches(0.8),
             "🏗️ 架构设计与实现亮点", font_size=36, color=WHITE, bold=True)

# Left: Architecture
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
             "系统架构", font_size=24, color=PRIMARY, bold=True)
arch_items = [
    "• 三层架构：Frontend → API → Database",
    "• 前端：React SPA + Nginx 反向代理",
    "• 后端：FastAPI 异步路由",
    "• ORM：SQLAlchemy 2.0 异步引擎",
    "• 认证：JWT Token + Bearer 拦截器",
    "• 部署：Docker Compose 一键编排",
]
add_bullet_list(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4),
                arch_items, font_size=18, color=DARK, spacing=Pt(8))

# Right: Highlights
add_text_box(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
             "实现亮点", font_size=24, color=PRIMARY, bold=True)
highlight_items = [
    "• 5 个新 Model + 5 个新 API 模块",
    "• 统一 CRUD 基类模式（代码复用）",
    "• 前后端分离，RESTful 风格接口",
    "• Zustand 轻量状态管理",
    "• 完整的错误处理和异常反馈",
    "• README 规范完善",
    "• 扩展性强，可继续添加新功能",
]
add_bullet_list(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4),
                highlight_items, font_size=18, color=DARK, spacing=Pt(8))

# DB Schema section
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), LIGHT_BG)
add_text_box(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(1.0),
             "数据库扩展表：favorites | addresses | reviews | coupons | user_coupons | logistics_tracking\n原有表不动，新增 6 张表，通过外键关联现有 users / products / orders 表",
             font_size=16, color=DARK)

# ========== Slide 10: Summary ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), PRIMARY)

add_text_box(slide, Inches(1.2), Inches(1.0), Inches(10), Inches(1.0),
             "🎯 总结", font_size=42, color=WHITE, bold=True)

summary_items = [
    "✅ 五大功能全部实现，显著提升电商平台完整性",
    "✅ 前后端协同开发，代码架构清晰可维护",
    "✅ 新增 5 个模型 + 6 张数据库表 + 12 个 API 端点",
    "✅ 前端 4 个新页面 + 2 个新组件 + 5 个新 API 模块",
    "✅ README 文档完善，Docker 部署就绪",
    "✅ 模块化设计，后续可继续扩展",
]
add_bullet_list(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(3.5),
                summary_items, font_size=22, color=WHITE, spacing=Pt(12))

# ========== Slide 11: Thank You ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.3), Inches(7.5), PRIMARY)

add_text_box(slide, Inches(1.2), Inches(2.0), Inches(10), Inches(1.5),
             "🙏 感谢聆听", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.2), Inches(3.8), Inches(10), Inches(1),
             "欢迎提问与交流", font_size=28, color=SECONDARY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.2), Inches(5.5), Inches(10), Inches(0.6),
             "EasyShop E-Commerce Platform · 期末考核项目演示", font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# Save
output_path = "C:/Users/14235/ecommerce-final/期末考核-EasyShop电商平台.pptx"
prs.save(output_path)
print(f"✅ PPT saved to: {output_path}")
